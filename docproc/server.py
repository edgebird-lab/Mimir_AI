"""docproc — isolated document text extractor (Zone-S-class: no network, no secrets, no host access
beyond a read-only /project/in). PDF/DOCX/PPTX/TXT parsers have an RCE history, so they run HERE,
far from the orchestrator's secrets. A malicious document that exploits a parser lands in a non-root,
no-network, no-secret container. Extracted text is returned to the orchestrator as UNTRUSTED data
(the caller taints it). Reachable only on the internal compose network, bearer-token gated.
"""
from __future__ import annotations

import os
from pathlib import Path

from starlette.applications import Starlette
from starlette.requests import Request
from starlette.responses import JSONResponse, PlainTextResponse
from starlette.routing import Route

import convert as conv                                  # aliased: the /convert handler shadows the name

IN_DIR = Path(os.environ.get("MIMIR_IN_DIR", "/project/in")).resolve()
OUT_DIR = Path(os.environ.get("MIMIR_OUT_DIR", "/project/out")).resolve()
PROJECT_DIR = Path(os.environ.get("MIMIR_PROJECT_DIR", "/project")).resolve()
TOKEN = os.environ.get("MIMIR_DOCPROC_TOKEN", "")
MAX_BYTES = int(os.environ.get("MIMIR_DOC_MAX_BYTES", str(80 * 1024 * 1024)))   # 80 MB file cap
MAX_TEXT = int(os.environ.get("MIMIR_DOC_MAX_TEXT", str(4 * 1024 * 1024)))      # 4 MB text cap
PSEUDO_PAGE = 2000                                                                # chars per pseudo-page
ALLOWED = {".pdf", ".docx", ".pptx", ".txt", ".md", ".markdown"}


def _authed(r: Request) -> bool:
    if not TOKEN:
        return True                                  # internal-net only; token optional but recommended
    import secrets as _s
    return _s.compare_digest(r.headers.get("authorization", "").removeprefix("Bearer ").strip(), TOKEN)


def _safe_path(rel: str) -> Path:
    p = (IN_DIR / rel.lstrip("/")).resolve()
    if IN_DIR not in p.parents and p != IN_DIR:
        raise PermissionError("path escapes /project/in")
    if not p.is_file():
        raise FileNotFoundError(rel)
    if p.stat().st_size > MAX_BYTES:
        raise ValueError(f"file too large (> {MAX_BYTES} bytes)")
    return p


def _safe_doc_path(rel: str) -> Path:
    """Path validation for /import: the viewer passes /project-relative paths (e.g. 'out/thesis.docx'
    or 'in/skript.pdf'), so resolve against the project root and require the result to stay inside the
    read-only doc dirs (/project/in or /project/out). Both are mounted read-only; neither holds secrets."""
    p = (PROJECT_DIR / rel.lstrip("/")).resolve()
    if not any(base == p or base in p.parents for base in (IN_DIR, OUT_DIR)):
        raise PermissionError("path escapes the project document directories")
    if not p.is_file():
        raise FileNotFoundError(rel)
    if p.stat().st_size > MAX_BYTES:
        raise ValueError(f"file too large (> {MAX_BYTES} bytes)")
    return p


def _blocks(text: str, size: int = PSEUDO_PAGE) -> list[str]:
    text = text.strip()
    return [text[i:i + size] for i in range(0, len(text), size)] or [""]


def _extract_pdf(p: Path) -> list[dict]:
    from pdfminer.high_level import extract_pages
    from pdfminer.layout import LTTextContainer
    pages = []
    for i, layout in enumerate(extract_pages(str(p))):
        txt = "".join(el.get_text() for el in layout if isinstance(el, LTTextContainer)).strip()
        pages.append({"n": i + 1, "text": txt})
        if sum(len(x["text"]) for x in pages) > MAX_TEXT:
            break
    return pages


def _extract_docx(p: Path) -> list[dict]:
    import docx
    d = docx.Document(str(p))
    full = "\n".join(par.text for par in d.paragraphs if par.text.strip())
    return [{"n": i + 1, "text": b} for i, b in enumerate(_blocks(full))]


def _extract_pptx(p: Path) -> list[dict]:
    from pptx import Presentation
    prs = Presentation(str(p))
    pages = []
    for i, slide in enumerate(prs.slides):
        parts = [sh.text for sh in slide.shapes if getattr(sh, "has_text_frame", False) and sh.text.strip()]
        pages.append({"n": i + 1, "text": "\n".join(parts).strip()})
    return pages


def _extract_text(p: Path) -> list[dict]:
    raw = p.read_text(encoding="utf-8-sig", errors="replace")[:MAX_TEXT]   # utf-8 + strip BOM (Windows-safe)
    return [{"n": i + 1, "text": b} for i, b in enumerate(_blocks(raw))]


EXTRACTORS = {".pdf": _extract_pdf, ".docx": _extract_docx, ".pptx": _extract_pptx,
              ".txt": _extract_text, ".md": _extract_text, ".markdown": _extract_text}


async def extract(request: Request):
    if not _authed(request):
        return PlainTextResponse("unauthorized", 401)
    body = await request.json()
    rel = str(body.get("path", ""))
    ext = os.path.splitext(rel)[1].lower()
    if ext not in ALLOWED:
        return JSONResponse({"error": f"unsupported type {ext}"}, 415)
    try:
        p = _safe_path(rel)
        pages = EXTRACTORS[ext](p)
    except (PermissionError, FileNotFoundError, ValueError) as e:
        return JSONResponse({"error": f"{type(e).__name__}: {e}"}, 400)
    except Exception as e:  # noqa: BLE001 — a parser blew up on a hostile file; contained here
        return JSONResponse({"error": f"parse failed: {type(e).__name__}: {e}"}, 422)
    pages = [pg for pg in pages if pg["text"]]
    chars = sum(len(pg["text"]) for pg in pages)
    return JSONResponse({"name": os.path.basename(rel), "ext": ext, "pages": pages,
                         "page_count": len(pages), "chars": chars})


async def convert(request: Request):
    """Canonical Markdown -> any target format (docx/pdf/html/odt/epub/pptx/rst/…) via the conversion hub.
    Returns base64. Isolated: the orchestrator sends the content and writes the returned bytes into out/."""
    if not _authed(request):
        return PlainTextResponse("unauthorized", 401)
    body = await request.json()
    content = str(body.get("content", ""))
    fmt = str(body.get("to", "docx")).lower()
    if fmt not in conv.EXPORT_FORMATS:
        return JSONResponse({"error": f"unsupported target {fmt}"}, 415)
    bib = body.get("bibliography")
    csl = str(body.get("csl", "apa"))
    # SVG math filter only on the pdf path (weasyprint can't render math); OMML handles docx natively.
    mf = "/app/mathfilter.py" if fmt == "pdf" and os.path.exists("/app/mathfilter.py") else None
    try:
        data = conv.from_markdown(content, fmt, bibliography=bib, csl=csl, math_filter=mf)
        return JSONResponse({"to": fmt, "data": conv.b64(data)})
    except ValueError as e:
        return JSONResponse({"error": str(e)}, 415)
    except Exception as e:  # noqa: BLE001 — a hostile document blew up a parser; contained here
        return JSONResponse({"error": f"{type(e).__name__}: {e}"}, 500)


async def import_doc(request: Request):
    """Any supported document -> canonical Markdown (the internal representation the viewer/corpus use).
    Reads only from /project/in; the returned Markdown is UNTRUSTED (the caller taints it)."""
    if not _authed(request):
        return PlainTextResponse("unauthorized", 401)
    body = await request.json()
    rel = str(body.get("path", ""))
    ext = os.path.splitext(rel)[1].lower()
    if ext not in conv.IMPORT_OK:
        return JSONResponse({"error": f"unsupported type {ext}"}, 415)
    try:
        p = _safe_doc_path(rel)
        result = conv.to_markdown(str(p), ext)
    except (PermissionError, FileNotFoundError, ValueError) as e:
        return JSONResponse({"error": f"{type(e).__name__}: {e}"}, 400)
    except Exception as e:  # noqa: BLE001 — parser blew up on a hostile file; contained here
        return JSONResponse({"error": f"import failed: {type(e).__name__}: {e}"}, 422)
    return JSONResponse({"name": os.path.basename(rel), "ext": ext, "markdown": result["markdown"],
                         "page_count": len(result.get("pages") or []), "chars": len(result["markdown"])})


async def health(request: Request):
    return JSONResponse({"ok": True, "supported": sorted(ALLOWED),
                         "import": sorted(conv.IMPORT_OK), "convert": sorted(conv.EXPORT_FORMATS)})


app = Starlette(routes=[
    Route("/extract", extract, methods=["POST"]),
    Route("/import", import_doc, methods=["POST"]),
    Route("/convert", convert, methods=["POST"]),
    Route("/health", health, methods=["GET"]),
])
