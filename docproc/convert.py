r"""Conversion hub — Markdown is the canonical intermediate format.

Everything dangerous (PDF/DOCX/HTML/EPUB parsers, which have an RCE history) runs HERE, inside the
isolated docproc container (no network, no secrets, non-root, read-only rootfs). Two directions go
through ONE hardened path:

  * to_markdown(path)          any supported format  ->  canonical Markdown  (import)
  * from_markdown(md, to, …)   canonical Markdown     ->  any target format   (export)

Single hardening point: every pandoc reader/writer that touches markup uses `-raw_html-raw_tex` so
injected <script>/<img …>/\input exfil markup can't survive a round-trip, while `+tex_math_dollars`
keeps $…$ / $$…$$ math (→ native Word OMML; the PDF path adds an SVG math filter). The Markdown canon
is inert (no active code), so it's a safe state to hand to the viewer or re-render.
"""
from __future__ import annotations

import base64
import json
import os
import subprocess
import tempfile

MAX_TEXT = int(os.environ.get("MIMIR_DOC_MAX_TEXT", str(4 * 1024 * 1024)))
PSEUDO_PAGE = 2000

# Scratch dir for pandoc/weasyprint. In the Linux container this is the /tmp tmpfs; on the native
# Windows build there is no /tmp, so fall back to the platform temp dir (cross-platform, same behaviour).
_TMP = "/tmp" if os.path.isdir("/tmp") else tempfile.gettempdir()


def _has_weasyprint() -> bool:
    """PDF export uses pandoc's --pdf-engine=weasyprint; the Windows build ships without the GTK stack
    weasyprint needs, so PDF is unavailable there while docx/html/odt/epub still work via pandoc."""
    import shutil
    return shutil.which("weasyprint") is not None

# Reader map: extension -> pandoc source format. PDF/PPTX are NOT here (pandoc can't read them) —
# they go through the dedicated python extractors below and are assembled into Markdown.
PANDOC_READ = {
    ".docx": "docx", ".odt": "odt", ".epub": "epub", ".rtf": "rtf",
    ".html": "html", ".htm": "html", ".tex": "latex", ".latex": "latex",
    ".rst": "rst", ".org": "org", ".textile": "textile", ".mediawiki": "mediawiki",
    ".md": "markdown", ".markdown": "markdown", ".csv": "csv", ".ipynb": "ipynb",
}
# Import formats served by dedicated python parsers (page-structured), not pandoc.
NATIVE_READ = {".pdf", ".pptx", ".txt"}
IMPORT_OK = set(PANDOC_READ) | NATIVE_READ

# Export targets (md -> X). weasyprint renders pdf; everything else is pandoc-native (no LibreOffice).
EXPORT_FORMATS = {"docx", "pdf", "html", "odt", "epub", "pptx", "rst", "gfm", "plain", "latex"}

_TMPENV = {**os.environ, "TMPDIR": _TMP, "HOME": _TMP}


def _reader_fmt(fmt: str) -> str:
    """Strip active content (raw HTML/TeX) on import for markup formats; keep math dollars."""
    if fmt in ("markdown", "html", "latex"):
        return f"{fmt}-raw_html-raw_tex+tex_math_dollars"
    return fmt


def _blocks(text: str, size: int = PSEUDO_PAGE) -> list[str]:
    text = text.strip()
    return [text[i:i + size] for i in range(0, len(text), size)] or [""]


def extract_pdf_pages(path: str) -> list[dict]:
    from pdfminer.high_level import extract_pages
    from pdfminer.layout import LTTextContainer
    pages, total = [], 0
    for i, layout in enumerate(extract_pages(path)):
        txt = "".join(el.get_text() for el in layout if isinstance(el, LTTextContainer)).strip()
        pages.append({"n": i + 1, "text": txt})
        total += len(txt)
        if total > MAX_TEXT:
            break
    return pages


def extract_pptx_pages(path: str) -> list[dict]:
    from pptx import Presentation
    pages = []
    for i, slide in enumerate(Presentation(path).slides):
        parts = [sh.text for sh in slide.shapes if getattr(sh, "has_text_frame", False) and sh.text.strip()]
        pages.append({"n": i + 1, "text": "\n".join(parts).strip()})
    return pages


def _pages_to_markdown(pages: list[dict], slides: bool = False) -> str:
    """Assemble page/slide text into Markdown with an HTML-comment page marker for provenance.
    The <!--page:n--> markers let corpus.py keep page references after import."""
    out = []
    for p in pages:
        if not p["text"]:
            continue
        out.append(f"<!--page:{p['n']}-->")
        if slides:
            out.append(f"## Folie {p['n']}\n")
        out.append(p["text"])
    return "\n\n".join(out).strip()


def to_markdown(path: str, ext: str) -> dict:
    """any supported document -> canonical Markdown. Returns {markdown, pages, ext}.
    `pages` is present for PDF/PPTX (page-structured) so RAG keeps provenance; empty otherwise."""
    ext = ext.lower()
    if ext == ".pdf":
        pages = extract_pdf_pages(path)
        return {"markdown": _pages_to_markdown(pages), "pages": pages, "ext": ext}
    if ext == ".pptx":
        pages = extract_pptx_pages(path)
        return {"markdown": _pages_to_markdown(pages, slides=True), "pages": pages, "ext": ext}
    if ext == ".txt":
        raw = open(path, encoding="utf-8", errors="replace").read()[:MAX_TEXT]
        pages = [{"n": i + 1, "text": b} for i, b in enumerate(_blocks(raw))]
        return {"markdown": raw, "pages": pages, "ext": ext}
    if ext in PANDOC_READ:
        src = _reader_fmt(PANDOC_READ[ext])
        p = subprocess.run(
            ["pandoc", "-f", src, "-t", "gfm", "--wrap=none", path],
            capture_output=True, timeout=180, cwd=_TMP, env=_TMPENV)
        if p.returncode != 0:
            raise RuntimeError("pandoc import: " + p.stderr.decode("utf-8", "replace")[:300])
        md = p.stdout.decode("utf-8", "replace")[:MAX_TEXT]
        return {"markdown": md, "pages": [], "ext": ext}
    raise ValueError(f"unsupported import type {ext}")


def from_markdown(content: str, to: str, bibliography=None, csl: str = "apa",
                  math_filter: str | None = None) -> bytes:
    """canonical Markdown -> target format bytes. weasyprint for pdf, pandoc-native otherwise.
    `bibliography` (CSL-JSON) + `csl` render citations; `math_filter` (a pandoc filter path) is applied
    for the pdf path so formulas become SVG that weasyprint can draw (it can't render MathML)."""
    to = to.lower()
    if to not in EXPORT_FORMATS:
        raise ValueError(f"unsupported target {to}")
    content = content[:MAX_TEXT]
    if to == "pdf" and not _has_weasyprint():
        raise ValueError("pdf export not available on this platform (needs weasyprint/GTK) - "
                         "export to docx/html/odt/epub instead")
    outp = tempfile.mktemp(suffix=f".{to}", dir=_TMP)
    cmd = ["pandoc", "--metadata", "title=Dokument", "--metadata", "lang=de-DE",
           "-f", "markdown-raw_html-raw_tex+tex_math_dollars", "-t", to,
           "--toc", "--toc-depth=3", "-o", outp]
    if to == "pdf":
        cmd += ["--pdf-engine=weasyprint"]
        if math_filter:
            cmd += ["--filter", math_filter]
    if bibliography:
        csl_path = os.path.join(os.environ.get("MIMIR_CSL_DIR", "/app/csl"), f"{os.path.basename(csl)}.csl")
        if os.path.exists(csl_path):
            bibp = tempfile.mktemp(suffix=".json", dir=_TMP)
            with open(bibp, "w") as bf:
                json.dump(bibliography, bf)
            cmd += ["--citeproc", "--bibliography=" + bibp, "--csl=" + csl_path]
    try:
        p = subprocess.run(cmd, input=content.encode("utf-8"), capture_output=True,
                           timeout=240, cwd="/tmp", env=_TMPENV)
        if p.returncode != 0 or not os.path.exists(outp):
            raise RuntimeError("pandoc: " + p.stderr.decode("utf-8", "replace")[:300])
        with open(outp, "rb") as f:
            return f.read()
    finally:
        try:
            os.remove(outp)
        except OSError:
            pass


def b64(data: bytes) -> str:
    return base64.b64encode(data).decode()
